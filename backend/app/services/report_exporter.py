"""简报导出器：Markdown → PPTX / HTML / TXT；BriefingDoc → PPTX / HTML"""
from __future__ import annotations

import base64
import io
import re
from datetime import datetime
from html import escape as _html_escape
from pathlib import Path
from typing import List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 默认主题色
THEME_PRIMARY = "#2563eb"   # 蓝
THEME_ACCENT  = "#10b981"   # 绿
THEME_DANGER  = "#ef4444"   # 红
THEME_WARNING = "#f59e0b"   # 橙
THEME_INFO    = "#3b82f6"   # 浅蓝
CHINESE_FONT  = "微软雅黑"

# matplotlib 字体兜底链(阶段 0.2 spike 验证可用)
MPL_FONT_CHAIN = ["Microsoft YaHei", "SimHei", "Noto Sans CJK SC", "PingFang SC"]


# ============================================================================
# Markdown 解析
# ============================================================================
def parse_markdown(text: str) -> List[Dict[str, Any]]:
    """
    把 markdown 切成章节
    第一个一级标题（# X）作为封面
    每个二级标题（## X）作为一页内容幻灯片
    """
    lines = text.split('\n')
    sections: List[Dict[str, Any]] = []
    current: Dict[str, Any] | None = None

    for line in lines:
        stripped = line.strip()
        # 一级标题
        m1 = re.match(r'^#\s+(.+)$', stripped)
        m2 = re.match(r'^##\s+(.+)$', stripped)
        if m1:
            if current:
                sections.append(current)
            current = {"level": 1, "title": m1.group(1).strip(), "body_lines": []}
        elif m2:
            if current:
                sections.append(current)
            current = {"level": 2, "title": m2.group(1).strip(), "body_lines": []}
        else:
            if current is not None:
                current["body_lines"].append(line)

    if current:
        sections.append(current)

    for s in sections:
        s["body"] = "\n".join(s["body_lines"]).strip()
        del s["body_lines"]

    return sections


# ============================================================================
# 颜色辅助
# ============================================================================
def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ============================================================================
# PPTX 导出
# ============================================================================
def md_to_pptx(
    markdown_text: str,
    out_path: Path | str,
    subtitle: str = "东风汽车 · 智汇车联数智简报",
) -> Path:
    """把 markdown 简报转成 PPTX"""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    sections = parse_markdown(markdown_text)
    if not sections:
        sections = [{"level": 1, "title": "决策简报", "body": markdown_text}]

    # 封面
    cover = sections[0]
    _add_cover(prs, cover["title"], subtitle)

    # 内容页：所有 level=2 的章节；如果没有 level=2，则用封面的 body 作为单页
    content_sections = [s for s in sections if s["level"] == 2]
    if not content_sections and cover.get("body"):
        _add_content_slide(prs, "正文", cover["body"])
    else:
        for sec in content_sections:
            _add_content_slide(prs, sec["title"], sec["body"])

    # 末页
    _add_end_slide(prs)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def _add_cover(prs: Presentation, title: str, subtitle: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(THEME_PRIMARY)
    bar.line.fill.background()

    # 底部细色条
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, prs.slide_height - Inches(0.3),
        prs.slide_width, Inches(0.3),
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = hex_to_rgb(THEME_ACCENT)
    bar2.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.6),
        prs.slide_width - Inches(1.6), Inches(1.5),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = CHINESE_FONT
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)

    # 副标题
    sub_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.6),
        prs.slide_width - Inches(1.6), Inches(0.6),
    )
    sp = sub_box.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.name = CHINESE_FONT
    sr.font.size = Pt(16)
    sr.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    # 日期
    date_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.4),
        prs.slide_width - Inches(1.6), Inches(0.5),
    )
    dp = date_box.text_frame.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = datetime.now().strftime('%Y 年 %m 月 %d 日')
    dr.font.name = CHINESE_FONT
    dr.font.size = Pt(14)
    dr.font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf)


def _add_content_slide(prs: Presentation, title: str, body: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 顶部细色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(THEME_PRIMARY)
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.45),
        prs.slide_width - Inches(1.2), Inches(0.7),
    )
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = CHINESE_FONT
    tr.font.size = Pt(28)
    tr.font.bold = True
    tr.font.color.rgb = hex_to_rgb(THEME_PRIMARY)

    # 标题下方分割线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25),
        Inches(0.6), Emu(20000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(THEME_ACCENT)
    line.line.fill.background()

    # 正文
    body_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5),
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(1.8),
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    body_lines = body.split('\n')
    first = True
    for raw in body_lines:
        ln = raw.rstrip()
        if not ln.strip():
            # 空行也保留作段落间距
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            continue

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        # 列表项
        is_list = bool(re.match(r'^\s*([-*]|\d+\.)\s+', ln))
        # 去掉列表符号
        if is_list:
            ln = re.sub(r'^\s*([-*]|\d+\.)\s+', '• ', ln)
            p.level = 0

        # 处理 markdown 加粗（简单的 **...**）
        parts = re.split(r'(\*\*[^*]+\*\*)', ln)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
            run.font.name = CHINESE_FONT
            run.font.size = Pt(16) if not is_list else Pt(15)
            run.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


def _add_end_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 居中色块
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(2.5),
        Inches(9.33), Inches(2.5),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = hex_to_rgb(THEME_PRIMARY)
    box.line.fill.background()

    txt = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.8),
        Inches(9.33), Inches(2.0),
    )
    tf = txt.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "感 谢 聆 听"
    r1.font.name = CHINESE_FONT
    r1.font.size = Pt(36)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "智汇车联 · 数智简报"
    r2.font.name = CHINESE_FONT
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xe5, 0xe7, 0xeb)


# ============================================================================
# HTML 导出（用于浏览器打印另存为 PDF，最稳跨平台）
# ============================================================================
def md_to_html(markdown_text: str, title: str = "决策简报") -> str:
    """把 markdown 转成可打印的 HTML"""
    body_html = _md_to_simple_html(markdown_text)
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8" />
<title>{title}</title>
<style>
  @page {{ size: A4; margin: 18mm 16mm; }}
  body {{
    font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
    color: #1f2937;
    line-height: 1.75;
    max-width: 800px;
    margin: 0 auto;
    padding: 16px;
  }}
  h1 {{ color: {THEME_PRIMARY}; border-bottom: 3px solid {THEME_PRIMARY}; padding-bottom: 8px; }}
  h2 {{ color: {THEME_PRIMARY}; margin-top: 32px; border-left: 4px solid {THEME_ACCENT}; padding-left: 12px; }}
  h3 {{ color: #374151; }}
  ul, ol {{ padding-left: 24px; }}
  li {{ margin: 6px 0; }}
  strong {{ color: #111827; }}
  blockquote {{
    border-left: 4px solid #cbd5e1;
    padding: 8px 16px;
    color: #4b5563;
    background: #f8fafc;
    margin: 16px 0;
  }}
  code {{ background: #f1f5f9; padding: 2px 6px; border-radius: 3px; font-size: 0.92em; }}
  pre {{ background: #0f172a; color: #e2e8f0; padding: 12px; border-radius: 6px; overflow: auto; }}
  hr {{ border: none; border-top: 1px solid #e5e7eb; margin: 24px 0; }}
  .footer {{ margin-top: 32px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #6b7280; font-size: 12px; }}
  @media print {{
    .no-print {{ display: none; }}
  }}
</style>
</head>
<body>
{body_html}
<div class="footer">
  东风汽车 · 智汇车联数智简报 · 生成于 {datetime.now().strftime('%Y-%m-%d %H:%M')}
</div>
<div class="no-print" style="position:fixed;top:8px;right:8px;">
  <button onclick="window.print()" style="padding:8px 16px;background:{THEME_PRIMARY};color:#fff;border:0;border-radius:4px;cursor:pointer;">打印 / 另存为 PDF</button>
</div>
</body>
</html>"""


def _md_to_simple_html(md: str) -> str:
    """非常简单的 markdown 转 HTML（够用版）"""
    lines = md.split('\n')
    out: List[str] = []
    in_ul = False
    in_ol = False
    in_pre = False

    def close_lists():
        nonlocal in_ul, in_ol
        if in_ul:
            out.append('</ul>')
            in_ul = False
        if in_ol:
            out.append('</ol>')
            in_ol = False

    for raw in lines:
        line = raw.rstrip()

        if line.strip().startswith('```'):
            close_lists()
            if in_pre:
                out.append('</pre>')
                in_pre = False
            else:
                out.append('<pre>')
                in_pre = True
            continue
        if in_pre:
            out.append(_escape_html(line))
            continue

        m_h = re.match(r'^(#{1,6})\s+(.+)$', line)
        if m_h:
            close_lists()
            level = len(m_h.group(1))
            out.append(f'<h{level}>{_inline_md(m_h.group(2))}</h{level}>')
            continue

        m_ul = re.match(r'^\s*[-*]\s+(.+)$', line)
        if m_ul:
            if not in_ul:
                close_lists()
                out.append('<ul>')
                in_ul = True
            out.append(f'<li>{_inline_md(m_ul.group(1))}</li>')
            continue

        m_ol = re.match(r'^\s*\d+\.\s+(.+)$', line)
        if m_ol:
            if not in_ol:
                close_lists()
                out.append('<ol>')
                in_ol = True
            out.append(f'<li>{_inline_md(m_ol.group(1))}</li>')
            continue

        if line.startswith('> '):
            close_lists()
            out.append(f'<blockquote>{_inline_md(line[2:])}</blockquote>')
            continue

        if not line.strip():
            close_lists()
            out.append('')
            continue

        close_lists()
        out.append(f'<p>{_inline_md(line)}</p>')

    close_lists()
    if in_pre:
        out.append('</pre>')
    return '\n'.join(out)


def _inline_md(text: str) -> str:
    """处理行内 markdown：**bold** *italic* `code`"""
    text = _escape_html(text)
    text = re.sub(r'\*\*([^*]+)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'\*([^*]+)\*', r'<em>\1</em>', text)
    text = re.sub(r'`([^`]+)`', r'<code>\1</code>', text)
    return text


def _escape_html(s: str) -> str:
    return (s.replace('&', '&amp;')
             .replace('<', '&lt;')
             .replace('>', '&gt;'))


# ============================================================================
# 阶段 4 新增：BriefingDoc → PPTX / HTML（结构化导出）
# ============================================================================
from .briefing_schema import (
    BriefingDoc, TrendSection, RankingSection, DistributionSection,
    AlertSection, TextSection,
)

_LEVEL_COLOR_HEX = {"high": THEME_DANGER, "warning": THEME_WARNING, "info": THEME_INFO}
_TONE_COLOR_HEX  = {"positive": THEME_ACCENT, "negative": THEME_DANGER, "neutral": THEME_PRIMARY}
_PIE_PALETTE = ["#2563eb", "#10b981", "#f59e0b", "#8b5cf6", "#ef4444", "#06b6d4", "#94a3b8"]


# ----------------------------------------------------------------------------
# 4.1 matplotlib chart 渲染(失败时调用方降级为表格)
# ----------------------------------------------------------------------------
def _setup_matplotlib_font() -> str | None:
    """配置 matplotlib 中文字体,返回成功使用的字体名;全部缺失返回 None。"""
    try:
        import matplotlib  # noqa
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm
        available = {f.name for f in fm.fontManager.ttflist}
        for name in MPL_FONT_CHAIN:
            if name in available:
                plt.rcParams["font.sans-serif"] = [name]
                plt.rcParams["axes.unicode_minus"] = False
                return name
    except Exception:
        pass
    return None


def render_trend_png(sec: TrendSection) -> bytes | None:
    """把 TrendSection 渲染为 PNG bytes;数据为空或字体缺失返回 None。"""
    if not sec.data:
        return None
    if not _setup_matplotlib_font():
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots(figsize=(8.5, 4.2), dpi=140)
        xs = [p.x for p in sec.data]
        ys = [p.y for p in sec.data]
        ax.plot(xs, ys, marker="o", linewidth=2.5, color=THEME_PRIMARY)
        ax.fill_between(xs, ys, alpha=0.15, color=THEME_PRIMARY)
        title = f"{sec.title}({sec.metric}{f' · {sec.unit}' if sec.unit else ''})"
        ax.set_title(title, fontsize=14, color="#1f2937", pad=10)
        ax.grid(True, alpha=0.3, linestyle="--")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        for x, y in zip(xs, ys):
            ax.annotate(f"{y}", (x, y), textcoords="offset points",
                        xytext=(0, 8), ha="center", fontsize=10, color=THEME_PRIMARY)
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format="png", facecolor="white", bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


def render_distribution_png(sec: DistributionSection) -> bytes | None:
    """把 DistributionSection 渲染为环形饼图 PNG bytes;失败返回 None。"""
    if not sec.data:
        return None
    if not _setup_matplotlib_font():
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots(figsize=(7, 4.5), dpi=140)
        labels = [p.label for p in sec.data]
        values = [p.value for p in sec.data]
        colors = [_PIE_PALETTE[i % len(_PIE_PALETTE)] for i in range(len(labels))]
        ax.pie(values, labels=labels, colors=colors, autopct="%1.1f%%",
               startangle=90, pctdistance=0.78, textprops={"fontsize": 11})
        centre = plt.Circle((0, 0), 0.55, fc="white")
        ax.add_artist(centre)
        ax.set_title(sec.title, fontsize=14, color="#1f2937", pad=10)
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format="png", facecolor="white", bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


# ----------------------------------------------------------------------------
# 4.2 doc_to_pptx
# ----------------------------------------------------------------------------
def doc_to_pptx(doc: BriefingDoc, out_path: Path | str) -> Path:
    """把 BriefingDoc 转成 PPTX:封面 + 摘要 + 各 section + 行动项 + 合规 + 末页。"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _add_doc_cover(prs, doc)
    if doc.executive_summary:
        _add_doc_summary_slide(prs, doc)
    for sec in doc.sections:
        if isinstance(sec, TrendSection):
            _add_trend_slide(prs, sec)
        elif isinstance(sec, RankingSection):
            _add_ranking_slide(prs, sec)
        elif isinstance(sec, DistributionSection):
            _add_distribution_slide(prs, sec)
        elif isinstance(sec, AlertSection):
            _add_alert_slide(prs, sec)
        elif isinstance(sec, TextSection):
            _add_content_slide(prs, sec.title, sec.body)
    if doc.actions:
        _add_actions_slide(prs, doc)
    _add_compliance_slide(prs, doc)
    _add_end_slide(prs)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def _add_doc_cover(prs: Presentation, doc: BriefingDoc) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 顶部主色块
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(THEME_PRIMARY); bar.line.fill.background()
    # 底部辅色细条
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        0, prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = hex_to_rgb(THEME_ACCENT); bar2.line.fill.background()

    # 周期 / 主题徽章
    badge = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6), Inches(0.4))
    bp = badge.text_frame.paragraphs[0]
    br = bp.add_run(); br.text = f"东风汽车 · 智汇车联 · {doc.meta.period or '决策简报'}"
    br.font.name = CHINESE_FONT; br.font.size = Pt(14); br.font.color.rgb = RGBColor(255, 255, 255)

    # 主标语(headline)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6),
        prs.slide_width - Inches(1.6), Inches(1.4))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = doc.cover.headline or doc.meta.title
    r.font.name = CHINESE_FONT; r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)

    # 副信息行
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2),
        prs.slide_width - Inches(1.6), Inches(0.5))
    sp = sub_box.text_frame.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = f"主题:{doc.meta.topic}  ·  生成于 {doc.meta.generated_at[:10]}"
    sr.font.name = CHINESE_FONT; sr.font.size = Pt(15); sr.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    # KPI 徽章条(横向最多 4 个)
    kpis = doc.cover.kpi_strip[:4]
    if kpis:
        n = len(kpis)
        total_w = prs.slide_width - Inches(1.6)
        cell_w = total_w / n - Inches(0.15)
        gap = Inches(0.15)
        start_x = Inches(0.8)
        y = Inches(5.0)
        for i, k in enumerate(kpis):
            x = start_x + (cell_w + gap) * i
            tone_color = hex_to_rgb(_TONE_COLOR_HEX.get(k.tone or "neutral", THEME_PRIMARY))
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, Inches(1.6))
            box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xf8, 0xfa, 0xfc)
            box.line.color.rgb = tone_color; box.line.width = Pt(1.5)
            tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), cell_w - Inches(0.3), Inches(1.3))
            tf2 = tb.text_frame; tf2.word_wrap = True
            p1 = tf2.paragraphs[0]; r1 = p1.add_run()
            r1.text = k.value; r1.font.name = CHINESE_FONT; r1.font.size = Pt(20); r1.font.bold = True
            r1.font.color.rgb = RGBColor(0x11, 0x18, 0x27)
            p2 = tf2.add_paragraph(); r2 = p2.add_run()
            r2.text = k.label; r2.font.name = CHINESE_FONT; r2.font.size = Pt(11)
            r2.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
            if k.delta:
                p3 = tf2.add_paragraph(); r3 = p3.add_run()
                r3.text = k.delta; r3.font.name = CHINESE_FONT; r3.font.size = Pt(11)
                r3.font.bold = True; r3.font.color.rgb = tone_color


def _add_doc_summary_slide(prs: Presentation, doc: BriefingDoc) -> None:
    slide = _new_blank_slide(prs, title="📌 摘要")
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
        prs.slide_width - Inches(1.2), prs.slide_height - Inches(2.0))
    tf = body_box.text_frame; tf.word_wrap = True
    sentences = re.split(r"(?<=[。!?！?])\s*", doc.executive_summary or "")
    sentences = [s.strip() for s in sentences if s.strip()][:4]
    labels = ["现状", "归因", "行动方向", "补充"]
    first = True
    for i, sent in enumerate(sentences):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r1 = p.add_run(); r1.text = f"【{labels[i] if i < len(labels) else '·'}】 "
        r1.font.name = CHINESE_FONT; r1.font.size = Pt(16); r1.font.bold = True
        r1.font.color.rgb = hex_to_rgb(THEME_PRIMARY)
        r2 = p.add_run(); r2.text = sent
        r2.font.name = CHINESE_FONT; r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


def _add_trend_slide(prs: Presentation, sec: TrendSection) -> None:
    slide = _new_blank_slide(prs, title=f"📈 {sec.title}")
    png = render_trend_png(sec)
    if png:
        slide.shapes.add_picture(io.BytesIO(png),
            Inches(0.8), Inches(1.5), width=Inches(11.7), height=Inches(4.6))
    else:
        _add_data_table_fallback(slide, ["X", "Y"], [[p.x, str(p.y)] for p in sec.data])
    if sec.insight:
        _add_insight_line(slide, sec.insight, top=Inches(6.3))


def _add_ranking_slide(prs: Presentation, sec: RankingSection) -> None:
    slide = _new_blank_slide(prs, title=f"🏆 {sec.title}")
    if sec.rows and sec.columns:
        _add_data_table(slide, sec.columns, [[str(c) for c in row] for row in sec.rows[:10]])
    else:
        _add_empty_hint(slide, "暂无排名数据")
    if sec.insight:
        _add_insight_line(slide, sec.insight, top=Inches(6.3))


def _add_distribution_slide(prs: Presentation, sec: DistributionSection) -> None:
    slide = _new_blank_slide(prs, title=f"🎯 {sec.title}")
    png = render_distribution_png(sec)
    if png:
        slide.shapes.add_picture(io.BytesIO(png),
            Inches(2.5), Inches(1.5), width=Inches(8), height=Inches(4.6))
    else:
        _add_data_table_fallback(slide, ["类别", "数值"],
            [[p.label, str(p.value)] for p in sec.data])
    if sec.insight:
        _add_insight_line(slide, sec.insight, top=Inches(6.3))


def _add_alert_slide(prs: Presentation, sec: AlertSection) -> None:
    color_hex = _LEVEL_COLOR_HEX.get(sec.level, THEME_WARNING)
    slide = _new_blank_slide(prs, title=f"⚠️ {sec.title}",
        bar_color_hex=color_hex)
    # 大色块 level 徽章
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.0), Inches(0.5), Inches(1.6), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = hex_to_rgb(color_hex)
    badge.line.fill.background()
    bf = badge.text_frame; bp = bf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = sec.level.upper()
    br.font.name = CHINESE_FONT; br.font.size = Pt(13); br.font.bold = True
    br.font.color.rgb = RGBColor(255, 255, 255)
    # msg
    msg_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6),
        prs.slide_width - Inches(1.6), Inches(2.5))
    tf = msg_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = sec.msg
    r.font.name = CHINESE_FONT; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)
    # 证据链(沿用 alert 旧 evidence_text 字符串列表)
    if sec.evidence_text:
        ev_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3),
            prs.slide_width - Inches(1.6), Inches(2.5))
        tf2 = ev_box.text_frame; tf2.word_wrap = True
        p1 = tf2.paragraphs[0]; r1 = p1.add_run(); r1.text = "证据链:"
        r1.font.name = CHINESE_FONT; r1.font.size = Pt(13); r1.font.bold = True
        r1.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
        for e in sec.evidence_text:
            p2 = tf2.add_paragraph(); r2 = p2.add_run(); r2.text = f"• {e}"
            r2.font.name = CHINESE_FONT; r2.font.size = Pt(14)
            r2.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


def _add_actions_slide(prs: Presentation, doc: BriefingDoc) -> None:
    slide = _new_blank_slide(prs, title="✅ 行动项")
    cols = ["优先级", "负责人", "行动", "截止"]
    rows = []
    for a in doc.actions:
        prio = {"high": "HIGH", "medium": "MID", "low": "LOW"}.get(a.priority, a.priority)
        rows.append([prio, a.owner, a.action, a.deadline])
    _add_data_table(slide, cols, rows, col_widths=[Inches(1.4), Inches(2.4), Inches(7.0), Inches(1.6)])


def _add_compliance_slide(prs: Presentation, doc: BriefingDoc) -> None:
    slide = _new_blank_slide(prs, title="🔒 合规与数据出处")
    # 大字统计
    stat_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
        prs.slide_width - Inches(1.6), Inches(1.5))
    tf = stat_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "已脱敏 "
    r1.font.name = CHINESE_FONT; r1.font.size = Pt(28); r1.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    r2 = p.add_run(); r2.text = f"{doc.compliance.masked_field_count} / {doc.compliance.total_field_count}"
    r2.font.name = CHINESE_FONT; r2.font.size = Pt(40); r2.font.bold = True
    r2.font.color.rgb = hex_to_rgb(THEME_PRIMARY)
    r3 = p.add_run(); r3.text = " 字段"
    r3.font.name = CHINESE_FONT; r3.font.size = Pt(28); r3.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    # 审计 ID
    if doc.meta.audit_id:
        au_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6),
            prs.slide_width - Inches(1.6), Inches(0.5))
        ap = au_box.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run(); ar.text = f"审计 ID:{doc.meta.audit_id}"
        ar.font.name = "Consolas"; ar.font.size = Pt(13); ar.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
    # findings 列表
    if doc.compliance.findings:
        fbox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5),
            prs.slide_width - Inches(1.6), Inches(2.5))
        tf2 = fbox.text_frame; tf2.word_wrap = True
        p0 = tf2.paragraphs[0]; r0 = p0.add_run(); r0.text = "已脱敏字段清单:"
        r0.font.name = CHINESE_FONT; r0.font.size = Pt(13); r0.font.bold = True
        r0.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
        for f in doc.compliance.findings:
            pf = tf2.add_paragraph(); rf = pf.add_run(); rf.text = f"• {f}"
            rf.font.name = CHINESE_FONT; rf.font.size = Pt(13)
            rf.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


# ---- 通用 PPTX helper ----
def _new_blank_slide(prs: Presentation, title: str, bar_color_hex: str = THEME_PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(bar_color_hex); bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4),
        prs.slide_width - Inches(1.2), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]; tr = tp.add_run(); tr.text = title
    tr.font.name = CHINESE_FONT; tr.font.size = Pt(26); tr.font.bold = True
    tr.font.color.rgb = hex_to_rgb(bar_color_hex)
    return slide


def _add_data_table(slide, columns: list, rows: list, col_widths=None) -> None:
    n_cols = len(columns); n_rows = len(rows) + 1
    if not col_widths:
        cw = (Inches(13.33) - Inches(1.2)) / n_cols
        col_widths = [cw] * n_cols
    total_w = sum(col_widths, Emu(0))
    left = (Inches(13.33) - total_w) / 2
    top = Inches(1.7)
    height = Inches(0.45) * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, height).table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    # 表头
    for j, c in enumerate(columns):
        cell = table.cell(0, j); cell.text = ""
        tf = cell.text_frame; p = tf.paragraphs[0]; r = p.add_run(); r.text = c
        r.font.name = CHINESE_FONT; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid(); cell.fill.fore_color.rgb = hex_to_rgb(THEME_PRIMARY)
    # 数据行
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i, j); cell.text = ""
            tf = cell.text_frame; p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            r.font.name = CHINESE_FONT; r.font.size = Pt(12)
            r.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)
            if i <= 3 and len(rows) >= 3:  # 前三名底色
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xfa, 0xfb, 0xff)


def _add_data_table_fallback(slide, cols, rows) -> None:
    """图表渲染失败时的兜底:展示数据表 + 提示"""
    hint = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
        prs.slide_width if False else Inches(11.7), Inches(0.4))
    hp = hint.text_frame.paragraphs[0]; hr = hp.add_run()
    hr.text = "(图表渲染失败,以下为原始数据)"
    hr.font.name = CHINESE_FONT; hr.font.size = Pt(11)
    hr.font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf); hr.font.italic = True
    if rows:
        _add_data_table(slide, cols, rows)


def _add_empty_hint(slide, msg: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.0))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = msg
    r.font.name = CHINESE_FONT; r.font.size = Pt(16); r.font.italic = True
    r.font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf)


def _add_insight_line(slide, text: str, top) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), top, Inches(11.7), Inches(0.7))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xf0, 0xf9, 0xff)
    box.line.color.rgb = hex_to_rgb(THEME_PRIMARY); box.line.width = Pt(0)
    tb = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.12), Inches(11.4), Inches(0.5))
    tf = tb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
    r.text = f"💡  {text}"
    r.font.name = CHINESE_FONT; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)


# ----------------------------------------------------------------------------
# 4.3 doc_to_html(A4 打印优化 + 卡片视觉)
# ----------------------------------------------------------------------------
def doc_to_html(doc: BriefingDoc) -> str:
    """把 BriefingDoc 转成 A4 打印友好的 HTML(图表用 base64 png 内嵌)。"""
    cover_html = _h_cover(doc)
    summary_html = _h_summary(doc) if doc.executive_summary else ""
    sections_html = "\n".join(_h_section(s) for s in doc.sections)
    actions_html = _h_actions(doc) if doc.actions else ""
    compliance_html = _h_compliance(doc)

    return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8" />
<title>{_html_escape(doc.meta.title)}</title>
<style>
@page {{ size: A4; margin: 14mm 12mm; }}
* {{ box-sizing: border-box; }}
body {{
  font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
  color: #1f2937;
  margin: 0;
  padding: 16px;
  background: #f9fafb;
}}
.briefing {{ max-width: 920px; margin: 0 auto; display: flex; flex-direction: column; gap: 18px; }}
.briefing > * {{ page-break-inside: avoid; }}
.cover {{
  border-radius: 18px;
  padding: 28px 30px;
  background: linear-gradient(135deg, #1e3a8a 0%, #2563eb 50%, #1d4ed8 100%);
  color: #fff;
  page-break-after: always;
}}
.cover .badge {{ display: inline-block; padding: 4px 12px; border-radius: 999px; background: rgba(255,255,255,0.2); font-size: 12px; font-weight: 600; margin-bottom: 12px; }}
.cover h1 {{ font-size: 30px; margin: 0 0 10px; line-height: 1.25; }}
.cover .subline {{ font-size: 13px; opacity: 0.85; margin-bottom: 18px; }}
.kpi-strip {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 10px; }}
.kpi-badge {{ background: rgba(255,255,255,0.65); border-radius: 12px; padding: 12px 14px; color: #111827; }}
.kpi-badge .kv {{ font-size: 20px; font-weight: 700; }}
.kpi-badge .kl {{ font-size: 11px; color: #6b7280; }}
.kpi-badge .kd {{ font-size: 12px; font-weight: 600; margin-top: 2px; }}
.kpi-positive {{ border-left: 3px solid {THEME_ACCENT}; }} .kpi-positive .kd {{ color: {THEME_ACCENT}; }}
.kpi-negative {{ border-left: 3px solid {THEME_DANGER}; }} .kpi-negative .kd {{ color: {THEME_DANGER}; }}
.kpi-neutral  {{ border-left: 3px solid {THEME_PRIMARY}; }} .kpi-neutral  .kd {{ color: {THEME_PRIMARY}; }}

.summary {{ background: #f8fafc; border-left: 4px solid {THEME_PRIMARY}; border-right: 1px solid #e5e7eb; border-top: 1px solid #e5e7eb; border-bottom: 1px solid #e5e7eb; border-radius: 14px; padding: 18px 22px; }}
.summary h3 {{ margin: 0 0 10px; font-size: 14px; color: #6b7280; }}
.summary p {{ margin: 6px 0; font-size: 14px; line-height: 1.7; }}
.summary .step {{ display: inline-block; background: {THEME_PRIMARY}; color: #fff; padding: 1px 8px; border-radius: 10px; font-size: 11px; font-weight: 600; margin-right: 8px; }}

.card {{ background: #fff; border: 1px solid #e5e7eb; border-radius: 14px; padding: 18px 20px; }}
.card h3 {{ margin: 0 0 10px; font-size: 16px; color: #111827; }}
.card .insight {{ margin-top: 10px; padding: 8px 12px; background: #f0f9ff; border-left: 3px solid {THEME_PRIMARY}; border-radius: 6px; font-size: 12px; }}

.alert {{ background: #fff; border: 1px solid #e5e7eb; border-left: 4px solid {THEME_WARNING}; border-radius: 14px; padding: 16px 20px; }}
.alert.high    {{ border-left-color: {THEME_DANGER};  background: #fef2f2; }}
.alert.warning {{ border-left-color: {THEME_WARNING}; background: #fffbeb; }}
.alert.info    {{ border-left-color: {THEME_INFO};    background: #eff6ff; }}
.alert h3 {{ margin: 0; font-size: 15px; }}
.alert .level {{ float: right; padding: 2px 10px; border-radius: 10px; font-size: 11px; font-weight: 700; background: rgba(0,0,0,0.06); }}
.alert .msg {{ margin: 8px 0 6px; font-size: 14px; font-weight: 500; }}
.alert .evidence {{ font-size: 12px; color: #4b5563; padding-left: 22px; margin: 6px 0 0; }}

table.rank {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
table.rank th {{ background: #f9fafb; padding: 8px 10px; text-align: left; color: #6b7280; font-size: 12px; font-weight: 600; border-bottom: 1px solid #e5e7eb; }}
table.rank td {{ padding: 8px 10px; border-bottom: 1px solid #f3f4f6; color: #1f2937; }}
table.rank tr.top td {{ background: #fafbff; font-weight: 500; }}
.rank-no {{ display: inline-flex; width: 22px; height: 22px; border-radius: 50%; background: #e5e7eb; color: #4b5563; align-items: center; justify-content: center; font-size: 11px; font-weight: 700; }}
.rank-no.gold {{ background: #fef3c7; color: #b45309; }}
.rank-no.silver {{ background: #e5e7eb; color: #4b5563; }}
.rank-no.bronze {{ background: #fed7aa; color: #b45309; }}

.chart-img {{ display: block; max-width: 100%; height: auto; margin: 0 auto; }}
.empty-chart {{ height: 180px; display: flex; align-items: center; justify-content: center; color: #9ca3af; background: #f9fafb; border-radius: 6px; font-size: 13px; }}

.actions {{ background: #fff; border: 1px solid #e5e7eb; border-radius: 14px; padding: 16px 18px; }}
.actions h3 {{ margin: 0 0 10px; font-size: 16px; }}
.action-row {{ display: grid; grid-template-columns: 60px 110px 1fr 100px; gap: 10px; padding: 8px 10px; border-radius: 8px; background: #f9fafb; border: 1px solid #e5e7eb; margin: 5px 0; font-size: 13px; align-items: center; }}
.prio {{ font-size: 11px; font-weight: 700; text-align: center; padding: 2px 6px; border-radius: 6px; }}
.prio.high   {{ background: #fee2e2; color: #b91c1c; }}
.prio.medium {{ background: #fef3c7; color: #b45309; }}
.prio.low    {{ background: #e0e7ff; color: #4338ca; }}
.action-row .deadline {{ font-family: ui-monospace, monospace; font-size: 12px; color: #6b7280; text-align: right; }}

.compliance {{ display: flex; flex-wrap: wrap; gap: 10px; padding: 12px 16px; border-radius: 12px; background: #f9fafb; border: 1px dashed #d1d5db; font-size: 12px; color: #4b5563; align-items: center; }}
.compliance .audit {{ font-family: ui-monospace, monospace; background: #fff; padding: 1px 8px; border-radius: 4px; color: {THEME_PRIMARY}; font-size: 11px; }}

@media print {{
  body {{ background: #fff; padding: 0; }}
  .no-print {{ display: none !important; }}
}}
.print-bar {{ position: fixed; top: 8px; right: 8px; }}
.print-bar button {{ padding: 8px 16px; background: {THEME_PRIMARY}; color: #fff; border: 0; border-radius: 6px; cursor: pointer; }}
</style>
</head>
<body>
<div class="briefing">
{cover_html}
{summary_html}
{sections_html}
{actions_html}
{compliance_html}
</div>
<div class="print-bar no-print"><button onclick="window.print()">打印 / 另存为 PDF</button></div>
</body>
</html>"""


def _h_cover(doc: BriefingDoc) -> str:
    period = _html_escape(doc.meta.period or "决策简报")
    headline = _html_escape(doc.cover.headline or doc.meta.title)
    sub = f"主题:{_html_escape(doc.meta.topic)}  ·  生成于 {_html_escape(doc.meta.generated_at[:10])}"
    if doc.meta.audit_id:
        sub += f"  ·  审计 {_html_escape(doc.meta.audit_id)}"
    kpi_html = ""
    if doc.cover.kpi_strip:
        items = []
        for k in doc.cover.kpi_strip:
            tone = k.tone or "neutral"
            delta = f'<div class="kd">{_html_escape(k.delta)}</div>' if k.delta else ""
            items.append(f'<div class="kpi-badge kpi-{tone}">'
                         f'<div class="kv">{_html_escape(k.value)}</div>'
                         f'<div class="kl">{_html_escape(k.label)}</div>'
                         f'{delta}</div>')
        kpi_html = f'<div class="kpi-strip">{"".join(items)}</div>'
    return (f'<section class="cover">'
            f'<div class="badge">东风汽车 · 智汇车联 · {period}</div>'
            f'<h1>{headline}</h1>'
            f'<div class="subline">{sub}</div>'
            f'{kpi_html}</section>')


def _h_summary(doc: BriefingDoc) -> str:
    sentences = [s.strip() for s in re.split(r"(?<=[。!?！?])\s*", doc.executive_summary or "") if s.strip()][:4]
    labels = ["现状", "归因", "行动方向", "补充"]
    body = "".join(
        f'<p><span class="step">{labels[i] if i < len(labels) else "·"}</span>{_html_escape(s)}</p>'
        for i, s in enumerate(sentences)
    )
    return f'<section class="summary"><h3>📌 摘要</h3>{body}</section>'


def _h_section(sec) -> str:
    if isinstance(sec, TrendSection):
        return _h_trend(sec)
    if isinstance(sec, RankingSection):
        return _h_ranking(sec)
    if isinstance(sec, DistributionSection):
        return _h_distribution(sec)
    if isinstance(sec, AlertSection):
        return _h_alert(sec)
    if isinstance(sec, TextSection):
        return f'<section class="card"><h3>{_html_escape(sec.title)}</h3><p>{_html_escape(sec.body)}</p></section>'
    return ""


def _h_trend(sec: TrendSection) -> str:
    png = render_trend_png(sec)
    if png:
        b64 = base64.b64encode(png).decode("ascii")
        chart_html = f'<img class="chart-img" src="data:image/png;base64,{b64}" alt="trend chart" />'
    else:
        chart_html = '<div class="empty-chart">暂无趋势数据</div>'
    delta_html = ""
    if sec.delta and sec.delta.value is not None:
        sign = "↑" if sec.delta.value > 0 else "↓" if sec.delta.value < 0 else "·"
        delta_html = f'<span style="float:right;font-size:12px;color:{THEME_PRIMARY};font-weight:600;">{sign} {_html_escape(sec.delta.baseline)} {abs(sec.delta.value):.1f}%</span>'
    insight = f'<div class="insight">💡 {_html_escape(sec.insight)}</div>' if sec.insight else ""
    return (f'<section class="card"><h3>📈 {_html_escape(sec.title)}{delta_html}</h3>'
            f'{chart_html}{insight}</section>')


def _h_ranking(sec: RankingSection) -> str:
    if not sec.rows or not sec.columns:
        body = '<div class="empty-chart">暂无排名数据</div>'
    else:
        thead = "".join(f'<th>{_html_escape(c)}</th>' for c in sec.columns)
        rows_html = []
        rank_class = ["gold", "silver", "bronze"]
        for i, row in enumerate(sec.rows[:15]):
            no_cls = rank_class[i] if i < 3 else ""
            tr_cls = ' class="top"' if i < 3 else ""
            cells = f'<td><span class="rank-no {no_cls}">{i+1}</span></td>'
            cells += "".join(f'<td>{_html_escape(str(c))}</td>' for c in row)
            rows_html.append(f"<tr{tr_cls}>{cells}</tr>")
        body = (f'<table class="rank"><thead><tr><th>#</th>{thead}</tr></thead>'
                f'<tbody>{"".join(rows_html)}</tbody></table>')
    insight = f'<div class="insight">💡 {_html_escape(sec.insight)}</div>' if sec.insight else ""
    return f'<section class="card"><h3>🏆 {_html_escape(sec.title)}</h3>{body}{insight}</section>'


def _h_distribution(sec: DistributionSection) -> str:
    png = render_distribution_png(sec)
    if png:
        b64 = base64.b64encode(png).decode("ascii")
        chart_html = f'<img class="chart-img" src="data:image/png;base64,{b64}" alt="distribution chart" />'
    else:
        chart_html = '<div class="empty-chart">暂无分布数据</div>'
    insight = f'<div class="insight">💡 {_html_escape(sec.insight)}</div>' if sec.insight else ""
    return f'<section class="card"><h3>🎯 {_html_escape(sec.title)}</h3>{chart_html}{insight}</section>'


def _h_alert(sec: AlertSection) -> str:
    level = sec.level or "warning"
    evidence = ""
    if sec.evidence_text:
        items = "".join(f"<li>{_html_escape(e)}</li>" for e in sec.evidence_text)
        evidence = f'<ul class="evidence">{items}</ul>'
    return (f'<section class="alert {level}">'
            f'<h3>⚠️ {_html_escape(sec.title)}<span class="level">{level.upper()}</span></h3>'
            f'<div class="msg">{_html_escape(sec.msg)}</div>{evidence}</section>')


def _h_actions(doc: BriefingDoc) -> str:
    rows = []
    for a in doc.actions:
        rows.append(f'<div class="action-row">'
                    f'<span class="prio {a.priority}">{a.priority.upper()}</span>'
                    f'<span><strong>{_html_escape(a.owner)}</strong></span>'
                    f'<span>{_html_escape(a.action)}</span>'
                    f'<span class="deadline">{_html_escape(a.deadline)}</span>'
                    f'</div>')
    return f'<section class="actions"><h3>✅ 行动项 · {len(doc.actions)} 项</h3>{"".join(rows)}</section>'


def _h_compliance(doc: BriefingDoc) -> str:
    audit = f'<span>审计 ID</span> <code class="audit">{_html_escape(doc.meta.audit_id)}</code>' if doc.meta.audit_id else ""
    findings = ""
    if doc.compliance.findings:
        items = "".join(f"<li>{_html_escape(f)}</li>" for f in doc.compliance.findings)
        findings = f'<details style="flex-basis:100%;margin-top:6px;"><summary>已脱敏字段清单 ({len(doc.compliance.findings)})</summary><ul style="margin:6px 0 0;padding-left:22px;font-size:12px;">{items}</ul></details>'
    return (f'<section class="compliance">'
            f'<span>🔒 合规过滤:已脱敏 <strong>{doc.compliance.masked_field_count}</strong> / '
            f'<strong>{doc.compliance.total_field_count}</strong> 字段</span>'
            f'<span style="margin-left:auto;">{audit}</span>'
            f'{findings}</section>')
